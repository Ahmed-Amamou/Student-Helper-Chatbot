import io

from pptx import Presentation


def parse_pptx(file_bytes: bytes) -> str:
    """Extract text from a PPTX file."""
    prs = Presentation(io.BytesIO(file_bytes))
    text_parts = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_text.append(text)
        if slide_text:
            text_parts.append("\n".join(slide_text))
    return "\n\n".join(text_parts)
